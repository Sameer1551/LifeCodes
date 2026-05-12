#!/usr/bin/env python3
"""
extract_texts.py
---------------
Pure‑Python bulk‑text extractor for common office, spreadsheet, PDF and web
documents.  No Java / Apache Tika required.

Enhancements:
* Multiprocessing support (--workers) for parallel extraction.
* Smart filtering of temporary Office files (~$*).
* Pandas configuration to prevent text truncation in cells.
* Robust handling of empty file lists.
"""

import argparse
import logging
import shutil
import subprocess
import sys
from concurrent.futures import ProcessPoolExecutor, as_completed
from pathlib import Path
from typing import List, Optional, Tuple

# --------------------------------------------------------------------------- #
# Dependencies Check
# --------------------------------------------------------------------------- #
# We check for critical dependencies here to provide a friendly error
# instead of a stack trace.
missing_deps = []
try:
    import pandas as pd
except ImportError:
    missing_deps.append("pandas")

try:
    from tqdm import tqdm
except ImportError:
    missing_deps.append("tqdm")

if missing_deps:
    print(
        f"Error: Missing required dependencies: {', '.join(missing_deps)}.\n"
        f"Please install them: pip install {' '.join(missing_deps)}"
    )
    sys.exit(1)

# --------------------------------------------------------------------------- #
# Global state
# --------------------------------------------------------------------------- #
_PANDOC_AVAILABLE: bool = shutil.which("pandoc") is not None

# Configure Pandas to show full content in cells (no truncation)
pd.set_option('display.max_colwidth', None)


# --------------------------------------------------------------------------- #
# Helper utilities
# --------------------------------------------------------------------------- #
def _clean_whitespace(text: str) -> str:
    """
    Normalise line endings, trim trailing spaces and collapse consecutive
    blank lines to a single empty line.
    """
    lines = [ln.rstrip() for ln in text.splitlines()]
    cleaned: List[str] = []
    prev_blank = False

    for ln in lines:
        if not ln:
            if not prev_blank:
                cleaned.append("")
            prev_blank = True
            continue
        cleaned.append(ln)
        prev_blank = False

    return "\n".join(cleaned).strip()


# --------------------------------------------------------------------------- #
# Extractor Functions
# --------------------------------------------------------------------------- #

def _extract_docx(fp: Path) -> str:
    import docx
    from docx.document import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    document = docx.Document(str(fp))
    out: List[str] = []

    def iter_block_items(parent):
        if isinstance(parent, Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            txt = block.text.strip()
            if not txt:
                continue
            style = getattr(block.style, "name", "")
            if style.lower().startswith("heading"):
                try:
                    level = int(style.split()[-1])
                    level = max(1, min(level, 6))
                except Exception:
                    level = 1
                out.append("#" * level + " " + txt)
            else:
                out.append(txt)
        elif isinstance(block, Table):
            for row in block.rows:
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                out.append(" | ".join(row_data))

    return "\n".join(out)


def _extract_pptx(fp: Path) -> str:
    import pptx
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pres = pptx.Presentation(str(fp))
    slides: List[str] = []

    def extract_shape_text(shape) -> List[str]:
        texts: List[str] = []
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            txt = shape.text.strip()
            if txt:
                texts.append(txt)
        if getattr(shape, "has_table", False) and shape.table:
            for row in shape.table.rows:
                row_data = [cell.text_frame.text.strip().replace("\n", " ") for cell in row.cells if cell.text_frame]
                if any(row_data):
                    texts.append(" | ".join(row_data))
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            # python-pptx < 0.6.18 compatibility check
            if hasattr(shape, 'shapes'):
                for child in shape.shapes:
                    texts.extend(extract_shape_text(child))
        return texts

    for i, slide in enumerate(pres.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            parts.extend(extract_shape_text(shape))
            
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("--- Speaker Notes ---\n" + notes)
                
        if parts:
            slide_body = "\n".join(parts)
            slides.append(f"--- Slide {i} ---\n{slide_body}")

    return "\n\n".join(slides)


def _extract_pdf(fp: Path) -> str:
    """
    Try PyMuPDF (fitz) first, then fallback to pdfminer.
    """
    out: List[str] = []
    
    # --- Fast Path: PyMuPDF ---
    try:
        import fitz  # type: ignore
        with fitz.open(str(fp)) as doc:
            # Check for encryption
            if doc.is_encrypted:
                # Try empty password
                if not doc.authenticate(""):
                    logging.warning("PDF %s is encrypted and could not be opened.", fp)
                    return f"[Encrypted PDF: {fp.name}]"
            
            for page_no in range(len(doc)):
                page = doc[page_no]
                txt = page.get_text("text").strip()
                if txt:
                    out.append(f"--- Page {page_no + 1} ---\n{txt}")
        
        if out:
            return "\n\n".join(out)
            
    except Exception as e:
        logging.debug("PyMuPDF failed for %s: %s – falling back to pdfminer", fp, e)

    # --- Fallback: pdfminer.six ---
    try:
        from pdfminer.high_level import extract_text as pdf_extract_text
        raw = pdf_extract_text(str(fp))
        pages = raw.split("\x0c")
        for i, pg in enumerate(pages, start=1):
            pg = pg.strip()
            if pg:
                out.append(f"--- Page {i} ---\n{pg}")
        return "\n\n".join(out)
    except Exception as e:
        logging.error("pdfminer also failed for %s: %s", fp, e)
        return f"[Failed to extract PDF: {fp.name}]"


def _extract_excel(fp: Path) -> str:
    xl = pd.ExcelFile(fp)
    out: List[str] = []

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, header=None, dtype=str)
        if df.empty:
            continue
        # Use to_string with max_colwidth=None (set globally) to avoid truncation
        txt = df.fillna("").to_string(index=False, header=False)
        out.append(f"=== Sheet: {sheet_name} ===\n{txt}")

    return "\n\n".join(out)


def _extract_csv(fp: Path) -> str:
    try:
        df = pd.read_csv(fp, header=None, dtype=str, sep=None, engine="python")
        if df.empty:
            return ""
        return df.fillna("").to_string(index=False, header=False)
    except Exception as e:
        logging.debug("Pandas CSV extraction failed for %s: %s - falling back to plain text", fp, e)
        return fp.read_text(encoding="utf-8", errors="ignore")


def _extract_html(fp: Path) -> str:
    from bs4 import BeautifulSoup
    html = fp.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "meta", "head", "title", "svg"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _extract_plain(fp: Path) -> str:
    return fp.read_text(encoding="utf-8", errors="ignore")


def _extract_with_pandoc(fp: Path) -> str:
    if not _PANDOC_AVAILABLE:
        raise RuntimeError("Pandoc not found. Install it to handle legacy formats (.doc, .ppt, .rtf).")
    try:
        result = subprocess.run(
            ["pandoc", "-f", "auto", "-t", "plain", str(fp)],
            capture_output=True,
            text=True,
            check=True,
        )
        return result.stdout
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Pandoc conversion failed: {e.stderr.strip()}")


# --------------------------------------------------------------------------- #
# Dispatcher
# --------------------------------------------------------------------------- #
def extract_text(fp: Path) -> str:
    ext = fp.suffix.lower()

    if ext == ".docx":
        return _extract_docx(fp)
    if ext == ".pptx":
        return _extract_pptx(fp)
    if ext == ".pdf":
        return _extract_pdf(fp)
    if ext in {".xlsx", ".xls"}:
        return _extract_excel(fp)
    if ext in {".csv", ".tsv"}:
        return _extract_csv(fp)
    if ext in {".html", ".htm"}:
        return _extract_html(fp)
    if ext in {".txt", ".md"}:
        return _extract_plain(fp)

    # Legacy binary Office files – try Pandoc
    if ext in {".doc", ".ppt", ".rtf"}:
        return _extract_with_pandoc(fp)

    return _extract_with_pandoc(fp)


# --------------------------------------------------------------------------- #
# I/O & Orchestration
# --------------------------------------------------------------------------- #
def process_single_file(fp: Path) -> Tuple[str, str]:
    """
    Worker function for multiprocessing.
    Returns (filename, extracted_text).
    """
    logging.debug("Processing %s", fp)
    text = extract_text(fp)
    return fp.name, _clean_whitespace(text)


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Bulk extract plain text from Office, PDF, CSV, HTML documents."
    )
    parser.add_argument(
        "inputs",
        nargs="+",
        help="Files or directories (scanned recursively).",
    )
    parser.add_argument(
        "-o", "--out-dir",
        default="extracted_texts",
        help="Directory for per-file outputs.",
    )
    parser.add_argument(
        "--combine",
        metavar="FILE",
        help="Merge all outputs into a single file.",
    )
    parser.add_argument(
        "--skip-unsupported",
        action="store_true",
        default=True,  # Default to skipping to avoid crashing on temp files
        help="Skip files that cannot be processed (default: True).",
    )
    parser.add_argument(
        "--fail-fast",
        action="store_true",
        help="Abort on first error.",
    )
    parser.add_argument(
        "-j", "--workers",
        type=int,
        default=4,
        help="Number of parallel workers (default: 4).",
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="Enable DEBUG logging.",
    )
    args = parser.parse_args()

    # Logging
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(asctime)s %(levelname)s %(message)s",
        datefmt="%H:%M:%S",
    )

    # Gather files
    all_files: List[Path] = []
    for raw in args.inputs:
        p = Path(raw).expanduser().resolve()
        if not p.exists():
            logging.warning("Path missing: %s", p)
            continue
        if p.is_file():
            all_files.append(p)
        elif p.is_dir():
            # Filter out temporary Office files (starting with ~$)
            all_files.extend([
                x for x in p.rglob("*") 
                if x.is_file() and not x.name.startswith("~$")
            ])

    if not all_files:
        logging.error("No files found.")
        sys.exit(1)

    logging.info("Found %d files. Processing with %d workers.", len(all_files), args.workers)

    # Create output directory
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    results: List[Tuple[str, str]] = []

    # Use ProcessPoolExecutor for CPU-bound extraction tasks
    # We use a try/finally block to ensure keyboard interrupts are handled
    try:
        with ProcessPoolExecutor(max_workers=args.workers) as executor:
            # Submit all tasks
            future_to_path = {
                executor.submit(process_single_file, fp): fp for fp in all_files
            }
            
            # Process as they complete
            for future in tqdm(as_completed(future_to_path), total=len(all_files), desc="Extracting"):
                fp = future_to_path[future]
                try:
                    filename, text = future.result()
                    results.append((filename, text))
                except Exception as exc:
                    msg = f"Failed {fp.name}: {exc}"
                    if args.fail_fast:
                        logging.error(msg)
                        executor.shutdown(wait=False, cancel_futures=True)
                        sys.exit(1)
                    else:
                        logging.warning(msg)
    except KeyboardInterrupt:
        logging.warning("Interrupted by user.")
        sys.exit(1)

    # Write outputs
    for filename, text in results:
        if not args.combine:
            out_path = out_dir / (Path(filename).stem + ".txt")
            out_path.write_text(text, encoding="utf-8")

    # Combine if requested
    if args.combine:
        combine_path = Path(args.combine).expanduser().resolve()
        combine_path.parent.mkdir(parents=True, exist_ok=True)
        with combine_path.open("w", encoding="utf-8") as f:
            for filename, text in results:
                f.write(f"===== {filename} =====\n")
                f.write(text)
                f.write("\n\n")
        logging.info("Combined corpus written to %s", combine_path)

    logging.info("Done. Files saved to '%s'", out_dir)


if __name__ == "__main__":
    main()
