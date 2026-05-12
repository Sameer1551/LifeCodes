#!/usr/bin/env python3
"""
extract_texts.py
---------------
Pure‑Python bulk‑text extractor for common office, spreadsheet, PDF and web
documents.  No Java / Apache Tika required.

Features
--------
* Handles .docx, .pptx, .pdf, .xlsx/.xls, .csv/.tsv, .html/.htm, .txt/.md
* Legacy .doc / .ppt / .rtf are supported via Pandoc (if installed)
* Fast PDF extraction with PyMuPDF (fitz) – falls back to pdfminer.six
* Recursively processes directories
* Optional `--combine` to merge all outputs into a single corpus file
* Progress bar (tqdm) + structured logging
* `--skip-unsupported` (default) or `--fail-fast` on extraction errors
* Clean, human‑readable output (slide/page markers, markdown headings, sheet titles)
"""

# Usage (examples):
#   python filestotext.py my_doc.docx
#   python filestotext.py my_folder/ --combine corpus.txt -v
#   python filestotext.py my_folder/ --fail-fast
#
# For full CLI help:
#   python filestotext.py --help

# --------------------------------------------------------------------------- #
# Imports – heavy optional libs are imported lazily inside their extractor
# --------------------------------------------------------------------------- #
import argparse
import logging
import shutil               # <-- moved up (fix #1)
import subprocess
import sys
from pathlib import Path
from typing import List

import pandas as pd
from tqdm import tqdm

# --------------------------------------------------------------------------- #
# Global state --------------------------------------------------------------- #
# --------------------------------------------------------------------------- #
# Detect Pandoc once at import time – we only need to know if the binary is
# reachable.  The flag is used by the legacy‑file extractor.
_PANDOC_AVAILABLE: bool = shutil.which("pandoc") is not None

# --------------------------------------------------------------------------- #
# Helper utilities
# --------------------------------------------------------------------------- #
def _clean_whitespace(text: str) -> str:
    """
    Normalise line endings, trim trailing spaces and collapse consecutive
    blank lines to a single empty line.
    """
    lines = [ln.rstrip() for ln in text.splitlines()]
    cleaned: List[str] = []
    prev_blank = False

    for ln in lines:
        if not ln:
            if not prev_blank:
                cleaned.append("")        # keep ONE empty line
            prev_blank = True
            continue
        cleaned.append(ln)
        prev_blank = False

    return "\n".join(cleaned).strip()


# --------------------------------------------------------------------------- #
# ① DOCX → plain text (preserve heading levels as markdown)
# --------------------------------------------------------------------------- #
def _extract_docx(fp: Path) -> str:
    import docx                     # lazy import
    from docx.document import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    document = docx.Document(str(fp))
    out: List[str] = []

    def iter_block_items(parent):
        if isinstance(parent, Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    for block in iter_block_items(document):
        if isinstance(block, Paragraph):
            txt = block.text.strip()
            if not txt:
                continue
            style = getattr(block.style, "name", "")
            if style.lower().startswith("heading"):
                # “Heading 2” → level 2 → ## heading
                try:
                    level = int(style.split()[-1])
                    level = max(1, min(level, 6))   # clamp to markdown range
                except Exception:
                    level = 1
                out.append("#" * level + " " + txt)
            else:
                out.append(txt)
        elif isinstance(block, Table):
            for row in block.rows:
                # Replace newlines within cells with spaces, then join with |
                row_data = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                out.append(" | ".join(row_data))

    return "\n".join(out)


# --------------------------------------------------------------------------- #
# ② PPTX → slide‑wise text
# --------------------------------------------------------------------------- #
def _extract_pptx(fp: Path) -> str:
    import pptx                     # lazy import
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    pres = pptx.Presentation(str(fp))
    slides: List[str] = []

    def extract_shape_text(shape) -> List[str]:
        texts: List[str] = []
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            txt = shape.text.strip()
            if txt:
                texts.append(txt)
        if getattr(shape, "has_table", False) and shape.table:
            for row in shape.table.rows:
                row_data = [cell.text_frame.text.strip().replace("\n", " ") for cell in row.cells if cell.text_frame]
                if any(row_data):
                    texts.append(" | ".join(row_data))
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                texts.extend(extract_shape_text(child))
        return texts

    for i, slide in enumerate(pres.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            parts.extend(extract_shape_text(shape))
            
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append("--- Speaker Notes ---\n" + notes)
                
        if parts:
            slide_body = "\n".join(parts)
            slides.append(f"--- Slide {i} ---\n{slide_body}")

    return "\n\n".join(slides)


# --------------------------------------------------------------------------- #
# ③ PDF → page‑wise text (fast PyMuPDF, fallback to pdfminer)
# --------------------------------------------------------------------------- #
def _extract_pdf(fp: Path) -> str:
    """
    Try to use PyMuPDF (fitz) first – it is usually ~10× faster than pdfminer.
    If the library is missing or the file cannot be read, fall back
    to pdfminer.six.
    """
    # ---------- Fast path – PyMuPDF ----------
    try:
        import fitz  # type: ignore
        with fitz.open(str(fp)) as doc:          # <-- context manager (fix #3)
            out: List[str] = []
            for page_no in range(len(doc)):
                page = doc[page_no]
                txt = page.get_text("text").strip()
                if txt:
                    out.append(f"--- Page {page_no + 1} ---\n{txt}")
            return "\n\n".join(out)
    except Exception as e:
        logging.debug("PyMuPDF failed for %s: %s – falling back to pdfminer", fp, e)

    # ---------- Fallback – pdfminer.six ----------
    from pdfminer.high_level import extract_text as pdf_extract_text
    raw = pdf_extract_text(str(fp))
    pages = raw.split("\x0c")            # pdfminer inserts form‑feed between pages
    out: List[str] = []
    for i, pg in enumerate(pages, start=1):
        pg = pg.strip()
        if pg:
            out.append(f"--- Page {i} ---\n{pg}")
    return "\n\n".join(out)


# --------------------------------------------------------------------------- #
# ④ EXCEL (xlsx / xls) → sheet‑wise tables
# --------------------------------------------------------------------------- #
def _extract_excel(fp: Path) -> str:
    xl = pd.ExcelFile(fp)
    out: List[str] = []

    for sheet_name in xl.sheet_names:
        df = xl.parse(sheet_name, header=None, dtype=str)
        if df.empty:
            continue
        txt = df.fillna("").to_string(index=False, header=False)
        out.append(f"=== Sheet: {sheet_name} ===\n{txt}")

    return "\n\n".join(out)


# --------------------------------------------------------------------------- #
# ⑤ CSV / TSV → plain table
# --------------------------------------------------------------------------- #
def _extract_csv(fp: Path) -> str:
    try:
        # pandas can infer the separator; we use the Python engine to avoid
        # C‑engine incompatibilities with `sep=None`.
        df = pd.read_csv(fp, header=None, dtype=str, sep=None, engine="python")
        if df.empty:
            return ""
        return df.fillna("").to_string(index=False, header=False)
    except Exception as e:
        logging.debug("Pandas CSV extraction failed for %s: %s - falling back to plain text", fp, e)
        return fp.read_text(encoding="utf-8", errors="ignore")


# --------------------------------------------------------------------------- #
# ⑥ HTML / HTM → visible text only
# --------------------------------------------------------------------------- #
def _extract_html(fp: Path) -> str:
    from bs4 import BeautifulSoup
    html = fp.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "meta", "head", "title", "svg"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


# --------------------------------------------------------------------------- #
# ⑦ Plain‑text fall‑backs (txt, md)
# --------------------------------------------------------------------------- #
def _extract_plain(fp: Path) -> str:
    return fp.read_text(encoding="utf-8", errors="ignore")


# --------------------------------------------------------------------------- #
# ⑧ LEGACY Office formats via Pandoc (if installed)
# --------------------------------------------------------------------------- #
def _extract_with_pandoc(fp: Path) -> str:
    """
    Uses the external `pandoc` binary:
        pandoc -f auto -t plain <file>
    Returns the plain‑text stdout.
    Raises RuntimeError on failure.
    """
    if not _PANDOC_AVAILABLE:
        raise RuntimeError(
            "Pandoc not found on the system. Install it to handle legacy "
            "formats (.doc, .ppt, .rtf)."
        )
    try:
        result = subprocess.run(
            ["pandoc", "-f", "auto", "-t", "plain", str(fp)],
            capture_output=True,
            text=True,
            check=True,
        )
        return result.stdout
    except subprocess.CalledProcessError as e:
        raise RuntimeError(f"Pandoc conversion failed: {e.stderr.strip()}")


# --------------------------------------------------------------------------- #
# Dispatcher – decide which extractor to call
# --------------------------------------------------------------------------- #
def extract_text(fp: Path) -> str:
    ext = fp.suffix.lower()

    if ext == ".docx":
        return _extract_docx(fp)
    if ext == ".pptx":
        return _extract_pptx(fp)
    if ext == ".pdf":
        return _extract_pdf(fp)
    if ext in {".xlsx", ".xls"}:
        return _extract_excel(fp)
    if ext in {".csv", ".tsv"}:
        return _extract_csv(fp)
    if ext in {".html", ".htm"}:
        return _extract_html(fp)
    if ext in {".txt", ".md"}:
        return _extract_plain(fp)

    # Legacy binary Office files – try Pandoc
    if ext in {".doc", ".ppt", ".rtf"}:
        return _extract_with_pandoc(fp)

    # Unknown extension – final Pandoc attempt (covers many exotic types)
    return _extract_with_pandoc(fp)


# --------------------------------------------------------------------------- #
# I/O helpers
# --------------------------------------------------------------------------- #
def _write_txt(txt: str, src: Path, out_dir: Path) -> Path:
    """Write cleaned text to `<out_dir>/<src‑stem>.txt` and return the created path."""
    out_path = out_dir / (src.stem + ".txt")
    out_path.write_text(_clean_whitespace(txt), encoding="utf-8")
    logging.debug("Wrote %s", out_path)
    return out_path


def process_one_file(src: Path, out_dir: Path) -> Path:
    """Extract text from *src* and store a .txt file in *out_dir*."""
    logging.info("Processing %s", src)
    raw = extract_text(src)
    return _write_txt(raw, src, out_dir)


# --------------------------------------------------------------------------- #
# CLI utilities
# --------------------------------------------------------------------------- #
def _gather_inputs(paths: List[str]) -> List[Path]:
    """Expand files / directories → flat list of Path objects."""
    all_files: List[Path] = []

    for raw in paths:
        p = Path(raw).expanduser().resolve()
        if not p.exists():
            logging.warning("Path does not exist – skipping: %s", p)
            continue

        if p.is_file():
            all_files.append(p)
        elif p.is_dir():
            # Recursively collect every file (filtering later)
            all_files.extend([x for x in p.rglob("*") if x.is_file()])
        else:
            logging.warning("Unsupported path type – skipping: %s", p)

    return all_files


def main() -> None:
    parser = argparse.ArgumentParser(
        description=(
            "Bulk extract plain text from Office, PDF, CSV, HTML and other documents "
            "without requiring Java/Tika."
        )
    )
    parser.add_argument(
        "inputs",
        nargs="+",
        help="One or more files and/or directories (directories are scanned recursively).",
    )
    parser.add_argument(
        "-o",
        "--out-dir",
        default="extracted_texts",
        help="Directory where per‑file .txt outputs will be stored (created if missing).",
    )
    parser.add_argument(
        "--combine",
        metavar="FILE",
        help=(
            "Write a single file that concatenates the output of all inputs. "
            "If omitted, each source gets its own .txt file."
        ),
    )
    parser.add_argument(
        "--skip-unsupported",
        action="store_true",
        help="Skip files that cannot be processed instead of aborting. "
             "(default behaviour – set this flag to explicitly enable it).",
    )
    parser.add_argument(
        "--fail-fast",
        action="store_true",
        help="Abort the whole run on the first extraction error.",
    )
    parser.add_argument(
        "--no-combine-header",
        action="store_true",
        help="When using --combine, omit the `===== filename =====` header before each file.",
    )
    parser.add_argument(
        "-v",
        "--verbose",
        action="store_true",
        help="Enable DEBUG‑level logging.",
    )
    args = parser.parse_args()

    # ------------------------------------------------------------------- #
    # Logging configuration
    # ------------------------------------------------------------------- #
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(asctime)s %(levelname)s %(message)s",
        datefmt="%H:%M:%S",
    )
    logging.debug("Pandoc available: %s", _PANDOC_AVAILABLE)

    # ------------------------------------------------------------------- #
    # Prepare output directory
    # ------------------------------------------------------------------- #
    out_dir = Path(args.out_dir).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    # ------------------------------------------------------------------- #
    # Gather all input files
    # ------------------------------------------------------------------- #
    src_files = _gather_inputs(args.inputs)
    if not src_files:
        logging.error("No readable files found – exiting.")
        sys.exit(1)

    logging.info("Found %d file(s) to process.", len(src_files))

    # ------------------------------------------------------------------- #
    # Main extraction loop
    # ------------------------------------------------------------------- #
    combined_chunks: List[str] = []

    for src in tqdm(src_files, desc="Extracting", unit="file"):
        try:
            txt_path = process_one_file(src, out_dir)

            if args.combine:
                header = "" if args.no_combine_header else f"===== {src.name} =====\n"
                combined_chunks.append(header + txt_path.read_text(encoding="utf-8"))
        except Exception as exc:
            msg = f"Failed to extract {src} – {exc}"
            if args.fail_fast:
                logging.error(msg)
                sys.exit(1)
            elif args.skip_unsupported:
                logging.warning(msg)
                continue
            else:
                logging.error(msg)
                # continue processing the rest (default behaviour)

    # ------------------------------------------------------------------- #
    # Write combined corpus if requested
    # ------------------------------------------------------------------- #
    if args.combine:
        combine_path = Path(args.combine).expanduser().resolve()
        combine_path.parent.mkdir(parents=True, exist_ok=True)
        combine_path.write_text("\n\n".join(combined_chunks), encoding="utf-8")
        logging.info("Combined corpus written to %s", combine_path)

    logging.info("All done – per‑file results live in '%s'", out_dir)


if __name__ == "__main__":
    main()
